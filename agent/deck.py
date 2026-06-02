"""Final deck compiler."""

from __future__ import annotations

from pathlib import Path

from agent.render import best_renderer, render_html_deck


def compile_deck_md(artefacts_dir: Path, context: dict) -> str:
    hackathon = context.get("hackathon_name", "EDTH Munich 2025")
    project = context.get("project_name", "Project X")
    team = context.get("team_name", "Team")

    def _slurp(name: str) -> str:
        p = artefacts_dir / name
        return p.read_text(encoding="utf-8") if p.exists() else f"[{name} not found]"

    problem = _slurp("02_candidate_problem.md")
    solution = _slurp("05_owner_pick.md")
    market = _slurp("07_market.md")
    comp = _slurp("07_competition.md")
    bm = _slurp("07_business_model.md")

    return "\n".join(
        [
            "---",
            "marp: true",
            "size: 16:9",
            "---",
            "",
            f"# {project}",
            f"**{hackathon}**",
            f"{team} | {context.get('date', '')}",
            "",
            "---",
            "",
            problem,
            "",
            "---",
            "",
            solution,
            "",
            "---",
            "",
            market,
            "",
            "---",
            "",
            comp,
            "",
            "---",
            "",
            bm,
            "",
            "---",
            "",
            "<!-- _class: lead -->",
            "",
            "# Thank You",
            f"**{team}**",
            "",
        ]
    )


def render_deck(artefacts_dir: Path, context: dict) -> Path:
    md_text = compile_deck_md(artefacts_dir, context)
    deck_md_path = artefacts_dir / "07_deck.md"
    deck_md_path.write_text(md_text, encoding="utf-8")

    renderer = best_renderer()
    if renderer == "marp":
        import subprocess

        from agent.render import detect_marp

        marp_path = detect_marp()
        try:
            subprocess.run(
                [str(marp_path), str(deck_md_path), "-o", str(artefacts_dir / "07_deck.html")],
                check=True,
                capture_output=True,
                timeout=60,
            )
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as exc:
            import logging

            logging.getLogger("agent.deck").warning(
                "Marp HTML render failed: %s. Falling back.", exc
            )
            return render_html_deck(artefacts_dir, md_text)
        from contextlib import suppress

        with suppress(subprocess.CalledProcessError, subprocess.TimeoutExpired):
            subprocess.run(
                [
                    str(marp_path),
                    str(deck_md_path),
                    "--pdf",
                    "-o",
                    str(artefacts_dir / "07_deck.pdf"),
                ],
                check=True,
                capture_output=True,
                timeout=60,
            )
        return artefacts_dir / "07_deck.html"
    if renderer == "pptx":
        from agent.render import render_pptx_deck

        return render_pptx_deck(artefacts_dir, md_text)
    return render_html_deck(artefacts_dir, md_text)


def render_pptx_deck(artefacts_dir: Path, md_text: str) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slides = md_text.split("\n---\n")
    for i, slide_md in enumerate(slides):
        if i == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            for line in slide_md.split("\n"):
                st = line.strip()
                if st.startswith("# "):
                    slide.shapes.title.text = st[2:]
                    break
                elif st.startswith("## "):
                    slide.shapes.title.text = st[3:]
                    break
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            lines = slide_md.strip().split("\n")
            if lines:
                first = lines[0].strip()
                slide.shapes.title.text = (
                    first.lstrip("#").strip()[:100] if first.startswith("#") else first[:100]
                )
    path = artefacts_dir / "07_deck.pptx"
    prs.save(str(path))
    return path
